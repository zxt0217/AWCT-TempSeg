#!/usr/bin/env python3
"""Replace implementation-style labels with paper-style functional labels.

This script edits only text in the current editable PPTX figure. It preserves
shape geometry, connectors, fills, borders, and the overall layout.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation


DEFAULT_INPUT = Path.home() / "Desktop" / "AWCT_TempSeg_Fig2_final_ready_submission.pptx"
DEFAULT_OUTPUT = Path.home() / "Desktop" / "AWCT_TempSeg_Fig2_short_labels_final.pptx"


REPLACEMENTS = {
    "Pair construction\nSemanticSTFTemporalDataset": ["Pair construction", "frame pairing"],
    "B1\nPT-v3m1": ["B1", "backbone"],
    "B2\nshared": ["B2", "backbone"],
    "B3\nshared": ["B3", "backbone"],
    "B4\nshared": ["B4", "backbone"],
    "TRF\nTemporal ref.": ["TRF", "reference"],
    "TRF\nFusion ref.": ["TRF", "fusion"],
    "Seg.\nLinear(C,19)": ["Seg.", "classifier"],
    "Semantic prediction\nseg_logits / labels": ["Semantic prediction", "point labels"],
    "shared t / t-1": ["shared frames"],
    "d=4": [""],
    "c=64": [""],
    "c=128": [""],
    "c=256": [""],
    "d=19": [""],
    "Global context pooling": ["Global reference"],
    "Local grid correspondence": ["Local reference"],
    "g_tm1_point": [""],
    "l_tm1_point, valid_match": [""],
    "Mini-batch sampling\nWeatherWeightedSampler": ["Mini-batch sampling", "weather-aware"],
    "VFB: Per-weather validation feedback": ["VFB: Weather validation feedback"],
    "Validation set\nSemanticSTF val": ["Validation set", "weather labels"],
    "Weather split\n_build_val_weather_map()": ["Weather split", "domain grouping"],
    "Weather-wise stats\nval mIoU / count": ["Weather-wise stats", "domain scores"],
    "Domain difficulty\nd = 1 - mIoU": ["Domain difficulty", "difficulty level"],
    "AWS: Adaptive weather curriculum sampler": ["AWS: Adaptive weather sampler"],
    "EMA\nmu=0.8": ["Smoothing", "stable feedback"],
    "Difficulty dist.\nsoftmax(D/tau)": ["Difficulty weighting", "weather balance"],
    "Curriculum prior\np_base": ["Curriculum prior", "initial balance"],
    "Conservative update\n(1-beta)p_base + beta q": ["Conservative update", "stable adjustment"],
    "Bounded projection\np_min <= p <= p_max": ["Bounded projection", "range control"],
    "Next-stage sampling\np_next -> sample weights": ["Next-stage sampling", "updated weights"],
    "Shared backbone block": ["Backbone block"],
    "Temporal reference/fusion module": ["Temporal reference"],
    "Per-weather validation feedback": ["Validation feedback"],
    "Adaptive weather sampler": ["Weather sampler"],
    "Feature stream": ["Feature flow"],
    "Curriculum feedback": ["Feedback flow"],
}


BANNED_TERMS = [
    "SemanticSTFTemporalDataset",
    "WeatherWeightedSampler",
    "_build_val_weather_map",
    "g_tm1_point",
    "l_tm1_point",
    "valid_match",
    "seg_logits",
    "Linear(C,19)",
    "val mIoU / count",
    "d = 1 - mIoU",
    "mu=0.8",
    "softmax",
    "tau",
    "p_base",
    "p_next",
    "p_min",
    "p_max",
    "beta",
    "sample weights",
    "<=",
    "PT-v3m1",
]


def set_text_preserve_style(shape, lines: list[str]) -> None:
    """Set paragraph text while preserving existing run style where possible."""
    tf = shape.text_frame
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()

    for idx, paragraph in enumerate(tf.paragraphs):
        value = lines[idx] if idx < len(lines) else ""
        if not paragraph.runs:
            paragraph.add_run()
        paragraph.runs[0].text = value
        for extra in paragraph.runs[1:]:
            extra.text = ""


def refine_labels(input_path: Path, output_path: Path) -> int:
    prs = Presentation(str(input_path))
    changed = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text") or not shape.text:
                continue
            current = shape.text
            if current in REPLACEMENTS:
                set_text_preserve_style(shape, REPLACEMENTS[current])
                changed += 1

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return changed


def parse_args():
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", default=str(DEFAULT_INPUT))
    parser.add_argument("--output", default=str(DEFAULT_OUTPUT))
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    input_path = Path(args.input).expanduser().resolve()
    output_path = Path(args.output).expanduser().resolve()
    changed = refine_labels(input_path, output_path)
    print(f"Wrote paper-style short-label PPTX to {output_path}")
    print(f"Updated text shapes: {changed}")


if __name__ == "__main__":
    main()
