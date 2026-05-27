#!/usr/bin/env python3
"""Draw Fig. 2 AWCT-TempSeg overall framework as editable PPT shapes.

This script intentionally uses python-pptx and editable PPT primitives
(rectangles, connectors, text boxes, small dot glyphs). It does not modify
training/model/config code.

Code-grounded module anchors:
- configs/semanticstf/semseg-pt-v3m1-0-tempseg-v2.py:
  TempSegV2Segmentor, PT-v3m1 backbone, SemanticSTFTemporalDataset.
- configs/semanticstf/semseg-pt-v3m1-0-tempseg-v2-awct-v11-conservative.py:
  ValidationGuidedWeatherCurriculumTrainer with tau/beta/EMA/bounds/p_base.
- pointcept/datasets/semanticstf.py:
  weather_map/_get_weather(), SemanticSTFTemporalDataset, prev_* fields, has_prev.
- pointcept/models/default.py:
  TempSegV2Segmentor, global/local temporal context, valid_match, temporal_fuse_*,
  residual feature fusion, seg_head.
- pointcept/engines/train.py:
  WeatherWeightedSampler, _build_val_weather_map(), weather-wise val mIoU,
  difficulty = 1 - mIoU, EMA, conservative update, bounded projection,
  p_next -> sample weights.

Conceptual blocks used for paper readability:
- "Temporal reference" groups _build_global_temporal_context() and
  _build_local_temporal_context().
- "Per-weather validation feedback" groups validation weather-map construction,
  weather-wise confusion matrices, val mIoU, and difficulty construction.
- "Adaptive weather curriculum sampler" groups the trainer-side update equations
  and WeatherWeightedSampler; it is not a single class with that exact name.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Arial"


class C:
    ink = "1F2328"
    muted = "5B6168"
    line = "2F3337"
    panel = "F8F9FB"
    panel_line = "C9D0D8"
    orange = "F8C9A1"
    orange2 = "FDE7D4"
    orange_line = "C85C12"
    blue = "CFE4F6"
    blue2 = "A8CBE8"
    blue_line = "2D77B8"
    green = "DCEED2"
    green2 = "B7D89E"
    green_line = "5F8E3E"
    purple = "E7DDF4"
    purple2 = "C2B6DF"
    purple_line = "6B55A3"
    gray = "ECEFF2"
    gray_line = "8E98A3"
    red_orange = "D96B1D"
    white = "FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def inches(v: float):
    return Inches(v)


def set_line(shape, color: str = C.line, width: float = 0.75, dash: bool = False) -> None:
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def set_fill(shape, color: str, transparency: float = 0.0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_text(
    shape,
    lines: list[tuple[str, float, bool, str | None]],
    *,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.04,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = inches(margin)
    tf.margin_right = inches(margin)
    tf.margin_top = inches(0.02)
    tf.margin_bottom = inches(0.02)
    for i, (text, size, bold, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color or C.ink)


def box(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[tuple[str, float, bool, str | None]],
    *,
    fill: str,
    line: str,
    radius: bool = False,
    dash: bool = False,
    lw: float = 0.75,
    name: str | None = None,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    if name:
        shape.name = name
    set_fill(shape, fill)
    set_line(shape, line, lw, dash)
    set_text(shape, lines)
    return shape


def textbox(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[tuple[str, float, bool, str | None]],
    *,
    align=PP_ALIGN.LEFT,
    name: str | None = None,
):
    shape = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    if name:
        shape.name = name
    set_text(shape, lines, align=align, margin=0.0)
    return shape


def panel(slide, x: float, y: float, w: float, h: float, title: str) -> None:
    shp = box(
        slide,
        x,
        y,
        w,
        h,
        [],
        fill=C.panel,
        line=C.panel_line,
        radius=True,
        dash=True,
        lw=0.7,
        name=f"Panel - {title}",
    )
    shp.fill.transparency = 8
    textbox(
        slide,
        x + 0.12,
        y + 0.06,
        w - 0.24,
        0.18,
        [(title, 7.6, True, C.muted)],
        name=f"Panel title - {title}",
    )


def add_arrowhead(connector) -> None:
    """python-pptx lacks arrowhead setters; add a DrawingML tailEnd node."""
    sp_pr = connector._element.spPr
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = OxmlElement("a:ln")
        sp_pr.append(ln)
    for old in ln.findall(qn("a:tailEnd")):
        ln.remove(old)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    ln.append(tail)


def arrow(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    *,
    color: str = C.line,
    width: float = 0.75,
    dash: bool = False,
    end: bool = True,
):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        inches(x1),
        inches(y1),
        inches(x2),
        inches(y2),
    )
    set_line(conn, color, width, dash)
    if end:
        add_arrowhead(conn)
    return conn


def dot(slide, x: float, y: float, color: str, size: float = 0.045) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, inches(x), inches(y), inches(size), inches(size))
    set_fill(shp, color)
    set_line(shp, color, 0.2)


def frame_icon(slide, x: float, y: float, w: float, h: float, title: str, subtitle: str, accent: str) -> None:
    box(
        slide,
        x,
        y,
        w,
        h,
        [(title, 7.6, True, C.ink), (subtitle, 5.8, False, C.muted)],
        fill=C.white,
        line=accent,
        radius=False,
        lw=0.75,
        name=title,
    )
    # Dot cloud glyph, deterministic and intentionally simple.
    pts = [
        (0.12, 0.22),
        (0.30, 0.17),
        (0.52, 0.25),
        (0.20, 0.43),
        (0.44, 0.48),
        (0.66, 0.39),
        (0.78, 0.56),
        (0.36, 0.66),
        (0.58, 0.70),
        (0.16, 0.72),
    ]
    for px, py in pts:
        dot(slide, x + 0.10 + px * (w - 0.20), y + 0.44 + py * (h - 0.55), accent, 0.035)


def output_icon(slide, x: float, y: float, w: float, h: float) -> None:
    box(
        slide,
        x,
        y,
        w,
        h,
        [],
        fill=C.green2,
        line=C.green_line,
        radius=False,
        lw=0.75,
        name="Semantic prediction output",
    )
    textbox(
        slide,
        x + 0.10,
        y + 0.13,
        w - 0.20,
        0.30,
        [("Semantic prediction", 7.3, True, C.ink), ("seg_logits / labels", 5.4, False, C.muted)],
        align=PP_ALIGN.CENTER,
    )
    colors = ["6AA84F", "3D85C6", "F1C232", "CC0000", "8E7CC3"]
    pts = [
        (0.18, 0.60, 0),
        (0.30, 0.72, 0),
        (0.42, 0.58, 1),
        (0.55, 0.74, 1),
        (0.70, 0.58, 2),
        (0.84, 0.69, 2),
        (0.48, 0.86, 3),
        (0.31, 0.86, 4),
        (0.73, 0.86, 4),
    ]
    for px, py, ci in pts:
        dot(slide, x + px * w, y + py * h, colors[ci], 0.05)


def legend_item(slide, x: float, y: float, color: str, label: str) -> None:
    box(slide, x, y + 0.03, 0.17, 0.11, [], fill=color, line=C.gray_line, radius=False, lw=0.45)
    textbox(slide, x + 0.22, y, 1.35, 0.18, [(label, 6.4, False, C.muted)])


def build_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C.white)

    textbox(
        slide,
        0.35,
        0.12,
        12.6,
        0.38,
        [("AWCT-TempSeg overall framework", 17.5, True, C.ink)],
        name="Title",
    )
    textbox(
        slide,
        0.36,
        0.47,
        12.6,
        0.18,
        [("Fig. 2-style editable framework diagram: temporal segmentation backbone with validation-guided adaptive weather curriculum", 6.7, False, C.muted)],
        name="Subtitle",
    )

    panel(slide, 0.30, 0.72, 12.75, 3.16, "AWCT-TempSeg main inference/training flow")
    panel(slide, 0.36, 4.08, 5.96, 1.76, "Per-weather validation feedback")
    panel(slide, 6.55, 4.08, 6.48, 1.76, "Adaptive weather curriculum sampler")

    # Main path input glyphs.
    frame_icon(slide, 0.62, 1.08, 1.30, 0.82, "Current frame", "coord + strength", C.orange_line)
    frame_icon(slide, 0.62, 2.04, 1.30, 0.82, "Historical frame", "prev_* fields", C.orange_line)
    box(
        slide,
        0.60,
        3.18,
        1.62,
        0.40,
        [("Mini-batch sampling", 7.0, True, C.ink), ("WeatherWeightedSampler", 5.8, False, C.muted)],
        fill=C.purple,
        line=C.purple_line,
        radius=False,
        lw=0.7,
        name="Mini-batch sampling",
    )
    arrow(slide, 1.28, 3.18, 1.24, 2.86, color=C.red_orange, width=0.7)
    arrow(slide, 1.39, 3.18, 1.24, 1.90, color=C.red_orange, width=0.7)

    # Shared backbone with two small internal lanes to avoid a business-card look.
    box(
        slide,
        2.35,
        1.28,
        2.10,
        1.35,
        [],
        fill=C.orange,
        line=C.orange_line,
        radius=False,
        lw=0.8,
        name="Shared backbone",
    )
    textbox(
        slide,
        2.50,
        1.40,
        1.80,
        0.30,
        [("Shared point cloud backbone", 8.0, True, C.ink), ("PT-v3m1 in TempSegV2Segmentor", 5.9, False, C.muted)],
        align=PP_ALIGN.CENTER,
    )
    box(slide, 2.57, 1.94, 0.72, 0.28, [("t", 6.2, True, C.orange_line), ("same weights", 4.9, False, C.muted)], fill=C.orange2, line=C.orange_line, radius=False, lw=0.45)
    box(slide, 3.49, 1.94, 0.72, 0.28, [("t-1", 6.2, True, C.orange_line), ("shared", 4.9, False, C.muted)], fill=C.orange2, line=C.orange_line, radius=False, lw=0.45)
    textbox(slide, 2.45, 2.73, 1.92, 0.18, [("paired temporal point clouds", 5.8, False, C.muted)], align=PP_ALIGN.CENTER)

    # Temporal reference module.
    box(
        slide,
        4.95,
        1.04,
        2.08,
        1.82,
        [("Temporal reference", 8.0, True, C.ink)],
        fill=C.blue,
        line=C.blue_line,
        radius=False,
        lw=0.8,
        name="Temporal reference",
    )
    box(slide, 5.17, 1.48, 1.64, 0.38, [("Global temporal context", 6.6, True, C.ink), ("frame-level mean pooling", 5.0, False, C.muted)], fill=C.white, line=C.blue_line, radius=False, lw=0.45)
    box(slide, 5.17, 2.05, 1.64, 0.42, [("Local same-grid correspondence", 6.4, True, C.ink), ("grid_coord hash / valid_match", 5.0, False, C.muted)], fill=C.white, line=C.blue_line, radius=False, lw=0.45)

    box(
        slide,
        7.55,
        1.15,
        1.95,
        1.60,
        [("Temporal feature fusion", 8.0, True, C.ink), ("temporal_fuse_global/local", 5.7, False, C.muted), ("has_prev + valid_match gates", 5.6, False, C.muted)],
        fill=C.blue2,
        line=C.blue_line,
        radius=False,
        lw=0.8,
        name="Temporal feature fusion",
    )
    box(
        slide,
        9.92,
        1.35,
        1.05,
        1.20,
        [("Segmentation", 7.3, True, C.ink), ("head", 7.3, True, C.ink), ("Linear(C, 19)", 5.6, False, C.muted)],
        fill=C.green,
        line=C.green_line,
        radius=False,
        lw=0.75,
        name="Segmentation head",
    )
    output_icon(slide, 11.42, 1.22, 1.34, 1.46)

    # Main arrows.
    arrow(slide, 1.92, 1.49, 2.35, 1.70)
    arrow(slide, 1.92, 2.45, 2.35, 2.18)
    arrow(slide, 4.45, 1.96, 4.95, 1.96)
    arrow(slide, 7.03, 1.96, 7.55, 1.96)
    arrow(slide, 9.50, 1.96, 9.92, 1.96)
    arrow(slide, 10.97, 1.96, 11.42, 1.96)

    # Tiny feature labels, in the style of paper pipeline annotations.
    textbox(slide, 2.36, 1.02, 2.06, 0.16, [("shared backbone for t and t-1", 5.7, False, C.muted)], align=PP_ALIGN.CENTER)
    textbox(slide, 5.03, 2.70, 1.92, 0.16, [("g_tm1_point + l_tm1_point", 5.7, False, C.muted)], align=PP_ALIGN.CENTER)
    textbox(slide, 7.60, 2.83, 1.84, 0.16, [("feat_out = feat_t + fused_global + fused_local", 5.4, False, C.muted)], align=PP_ALIGN.CENTER)

    # Validation feedback branch.
    y = 4.68
    box(slide, 0.66, y, 1.10, 0.58, [("Validation set", 6.8, True, C.ink), ("SemanticSTF val", 5.2, False, C.muted)], fill=C.blue, line=C.blue_line, radius=False, lw=0.65)
    box(slide, 2.05, y, 1.18, 0.58, [("Split by weather", 6.6, True, C.ink), ("_build_val_weather_map()", 5.0, False, C.muted)], fill=C.blue, line=C.blue_line, radius=False, lw=0.65)
    box(slide, 3.52, y - 0.10, 1.28, 0.78, [("Weather-wise statistics", 6.4, True, C.ink), ("CM / val mIoU", 5.0, False, C.muted), ("sample + point counts", 5.0, False, C.muted)], fill=C.blue, line=C.blue_line, radius=False, lw=0.65)
    box(slide, 5.08, y - 0.02, 0.98, 0.64, [("Weather-domain", 6.2, True, C.ink), ("difficulty", 6.2, True, C.ink), ("d = 1 - mIoU", 5.0, False, C.muted)], fill=C.purple, line=C.purple_line, radius=False, lw=0.65)
    arrow(slide, 1.76, y + 0.29, 2.05, y + 0.29)
    arrow(slide, 3.23, y + 0.29, 3.52, y + 0.29)
    arrow(slide, 4.80, y + 0.29, 5.08, y + 0.29)
    # From prediction to validation branch, subtle dashed evaluation path.
    arrow(slide, 12.10, 2.68, 12.10, 3.68, color=C.green_line, width=0.65, dash=True, end=False)
    arrow(slide, 12.10, 3.68, 1.20, 3.68, color=C.green_line, width=0.65, dash=True, end=False)
    arrow(slide, 1.20, 3.68, 1.20, 4.68, color=C.green_line, width=0.65, dash=True)

    # Adaptive weather curriculum sampler branch.
    arrow(slide, 6.06, y + 0.30, 6.76, y + 0.30, color=C.purple_line, width=0.75)
    box(slide, 6.76, y - 0.02, 0.92, 0.64, [("EMA", 6.6, True, C.ink), ("mu = 0.8", 5.2, False, C.muted)], fill=C.purple, line=C.purple_line, radius=False, lw=0.65)
    box(slide, 8.02, y - 0.02, 1.04, 0.64, [("Difficulty", 6.3, True, C.ink), ("distribution", 6.3, True, C.ink), ("softmax(D/tau)", 5.0, False, C.muted)], fill=C.purple, line=C.purple_line, radius=False, lw=0.65)
    box(slide, 9.42, y - 0.02, 1.08, 0.64, [("Fixed prior", 6.4, True, C.ink), ("p_base ratios", 5.0, False, C.muted)], fill=C.gray, line=C.gray_line, radius=False, lw=0.65)
    box(slide, 7.02, 5.38, 1.30, 0.46, [("Conservative update", 6.0, True, C.ink), ("(1-beta)p_base + beta q", 4.9, False, C.muted)], fill=C.purple, line=C.purple_line, radius=False, lw=0.65)
    box(slide, 8.70, 5.38, 1.22, 0.46, [("Bounded projection", 6.0, True, C.ink), ("0.10 <= p <= 0.45", 4.9, False, C.muted)], fill=C.purple, line=C.purple_line, radius=False, lw=0.65)
    box(slide, 10.32, 5.30, 1.92, 0.62, [("Next-stage sampling distribution", 6.2, True, C.ink), ("p_next -> sample weights", 5.0, False, C.muted)], fill=C.purple2, line=C.purple_line, radius=False, lw=0.65)
    arrow(slide, 7.68, y + 0.30, 8.02, y + 0.30)
    arrow(slide, 9.06, y + 0.30, 9.42, y + 0.30)
    arrow(slide, 8.54, y + 0.62, 7.70, 5.38)
    arrow(slide, 9.96, y + 0.62, 8.20, 5.38)
    arrow(slide, 8.32, 5.61, 8.70, 5.61)
    arrow(slide, 9.92, 5.61, 10.32, 5.61)

    # Dashed feedback loop: p_next updates the next epoch sampler.
    arrow(slide, 11.28, 5.92, 11.28, 6.08, color=C.red_orange, width=0.7, dash=True, end=False)
    arrow(slide, 11.28, 6.08, 0.45, 6.08, color=C.red_orange, width=0.7, dash=True, end=False)
    arrow(slide, 0.45, 6.08, 0.45, 3.38, color=C.red_orange, width=0.7, dash=True, end=False)
    arrow(slide, 0.45, 3.38, 0.60, 3.38, color=C.red_orange, width=0.7, dash=True)
    textbox(slide, 8.42, 6.13, 2.55, 0.16, [("next epoch sampling update", 5.6, False, C.red_orange)], align=PP_ALIGN.CENTER)

    # Legend and caption.
    textbox(slide, 3.15, 6.35, 0.55, 0.18, [("Legend", 6.4, True, C.muted)])
    legend_item(slide, 3.75, 6.34, C.orange, "backbone / input")
    legend_item(slide, 5.30, 6.34, C.blue2, "temporal reference / fusion")
    legend_item(slide, 7.30, 6.34, C.green2, "prediction")
    legend_item(slide, 8.72, 6.34, C.purple2, "feedback / sampler")
    textbox(
        slide,
        0.35,
        6.95,
        12.63,
        0.28,
        [("Fig. 2. Overall framework of the proposed AWCT-TempSeg.", 9.0, False, C.ink)],
        align=PP_ALIGN.CENTER,
        name="Caption",
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate editable Fig. 2 AWCT-TempSeg framework PPTX.")
    parser.add_argument(
        "--output",
        default=str(Path.home() / "Desktop" / "AWCT_TempSeg_Fig2_Framework.pptx"),
        help="Output pptx path.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    prs = Presentation()
    prs.slide_width = inches(SLIDE_W)
    prs.slide_height = inches(SLIDE_H)
    build_slide(prs)
    output = Path(args.output).expanduser().resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Wrote editable python-pptx figure to {output}")


if __name__ == "__main__":
    main()
