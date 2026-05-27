#!/usr/bin/env python3
"""SCI-style refinement for the current AWCT-TempSeg Fig. 2 PPT.

The logical structure is kept from tools/redraw_awct_fig2_shape_ppt.py:
Temporal LiDAR pair -> M0 -> B/TRF trunk -> F1/F2 -> Seg. -> prediction,
with VFB and AWS mechanisms below. This version only tightens typography,
spacing, line weights, feedback paths, and legend/caption placement.

All objects are editable PowerPoint shapes/connectors/text boxes. No PNG or
screenshot is inserted.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Arial"
CAPTION_FONT = "Times New Roman"


class C:
    ink = "111111"
    dim = "555555"
    line = "333333"
    skip = "B8B8B8"
    orange = "F6C7A4"
    orange_dark = "BD5A13"
    blue = "B7D1EC"
    blue_light = "DDECF8"
    blue_dark = "2F75B5"
    green = "BFDCA9"
    green_light = "E0EFD8"
    green_dark = "5A8B3F"
    purple = "CCC0E3"
    purple_light = "EDE6F6"
    purple_dark = "6A56A1"
    gray = "F7F7F7"
    gray_line = "9E9E9E"
    white = "FFFFFF"
    feedback = "C97936"
    val = "9BCB8F"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def I(v: float):
    return Inches(v)


def set_text(shape, lines, *, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.02):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(0.012)
    tf.margin_bottom = I(0.012)
    for i, (value, size, bold, color, *font_name) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = value
        run.font.name = font_name[0] if font_name else FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color or C.ink)


def box(slide, x, y, w, h, lines, fill, line=C.line, lw=1.0, name=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if name:
        shp.name = name
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(lw)
    set_text(shp, lines)
    return shp


def text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, name=None):
    shp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    if name:
        shp.name = name
    set_text(shp, lines, align=align, margin=0)
    return shp


def add_arrowhead(connector):
    sp_pr = connector._element.spPr
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = OxmlElement("a:ln")
        sp_pr.append(ln)
    for old in ln.findall(qn("a:tailEnd")):
        ln.remove(old)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    ln.append(tail)


def arrow(slide, x1, y1, x2, y2, color=C.line, lw=0.85, dash=False, head=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    c.line.color.rgb = rgb(color)
    c.line.width = Pt(lw)
    if dash:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if head:
        add_arrowhead(c)
    return c


def dot(slide, x, y, color, size=0.020):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(size), I(size))
    d.fill.solid()
    d.fill.fore_color.rgb = rgb(color)
    d.line.color.rgb = rgb(color)
    d.line.width = Pt(0.1)
    return d


def stack_block(slide, x, y, w, h, label, sub, fill, line):
    # Minimal stacked offset only, to echo E/D blocks in the reference without
    # creating a PPT shadow.
    for dx in (0.035, 0.070):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x + dx), I(y + dx), I(w), I(h))
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(0.45)
    return box(slide, x, y, w, h, [(label, 8.2, True, C.ink), (sub, 5.5, False, C.dim)], fill, line, 1.0)


def lidar_frame(slide, x, y, title):
    box(slide, x, y, 1.12, 0.50, [], C.white, C.orange_dark, 0.9)
    text(slide, x + 0.06, y + 0.04, 1.00, 0.15, [(title, 6.3, True, C.ink)], PP_ALIGN.CENTER)
    pts = [
        (0.11, 0.27), (0.18, 0.34), (0.28, 0.24), (0.39, 0.37),
        (0.51, 0.30), (0.63, 0.22), (0.76, 0.37), (0.88, 0.29),
        (0.98, 0.39), (0.42, 0.43), (0.23, 0.43),
    ]
    for px, py in pts:
        dot(slide, x + px, y + py, C.orange_dark, 0.020)
    for k in range(2):
        arrow(slide, x + 0.10, y + 0.42 - 0.07 * k, x + 0.98, y + 0.42 - 0.07 * k, "C4C4C4", 0.3, head=False)


def semantic_output(slide, x, y):
    box(slide, x, y, 1.25, 0.68, [], C.white, C.green_dark, 0.9)
    text(slide, x + 0.07, y + 0.04, 1.11, 0.20, [("Semantic prediction", 6.2, True, C.ink), ("seg_logits / labels", 5.1, False, C.dim)], PP_ALIGN.CENTER)
    colors = ["7AB36A", "4A8BC2", "E3B83B", "C63D32", "9788C8", "35A7C9", "E8A6C5"]
    pts = [
        (0.17, 0.46, 0), (0.29, 0.60, 0), (0.43, 0.47, 1),
        (0.55, 0.64, 1), (0.68, 0.46, 2), (0.81, 0.60, 2),
        (0.48, 0.80, 3), (0.31, 0.78, 4), (0.73, 0.77, 4),
        (0.58, 0.48, 5), (0.90, 0.72, 6),
    ]
    for px, py, idx in pts:
        dot(slide, x + px * 1.25, y + py * 0.68, colors[idx], 0.028)


def skip(slide, x0, x1, y_top, y_base):
    arrow(slide, x0, y_base, x0, y_top, C.skip, 0.65, head=False)
    arrow(slide, x0, y_top, x1, y_top, C.skip, 0.65, head=False)
    arrow(slide, x1, y_top, x1, y_base, C.skip, 0.65, head=True)


def legend(slide):
    x, y = 6.45, 5.72
    text(slide, x, y, 0.48, 0.15, [("Legend", 6.4, True, C.ink)])
    items = [
        ("B", "Shared backbone block", C.orange),
        ("TRF", "Temporal reference/fusion module", C.blue),
        ("VFB", "Per-weather validation feedback", C.blue_light),
        ("AWS", "Adaptive weather sampler", C.purple),
    ]
    for i, (code, desc, fill) in enumerate(items):
        yy = y + 0.22 * (i + 1)
        box(slide, x + 0.62, yy, 0.36, 0.15, [(code, 5.2, True, C.ink)], fill, C.gray_line, 0.55)
        text(slide, x + 1.04, yy - 0.006, 1.85, 0.15, [(desc, 5.2, False, C.dim)])
    arrow(slide, x + 3.10, y + 0.27, x + 3.58, y + 0.27, C.line, 0.75)
    text(slide, x + 3.66, y + 0.205, 1.00, 0.15, [("Feature stream", 5.2, False, C.dim)])
    arrow(slide, x + 3.10, y + 0.58, x + 3.58, y + 0.58, C.feedback, 0.65, dash=True)
    text(slide, x + 3.66, y + 0.515, 1.35, 0.15, [("Curriculum feedback", 5.2, False, C.dim)])


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C.white)

    # No large PPT title: keep only a tiny in-figure label.
    text(slide, 0.45, 0.20, 1.65, 0.14, [("AWCT-TempSeg", 6.5, True, C.dim)])

    # Top trunk: compact, paper-like network chain.
    lidar_frame(slide, 0.40, 0.98, "Current frame t")
    lidar_frame(slide, 0.40, 1.58, "Historical frame t-1")
    text(slide, 0.38, 0.76, 1.18, 0.15, [("Temporal LiDAR pair", 6.7, True, C.ink)], PP_ALIGN.CENTER)

    box(slide, 1.78, 1.34, 0.50, 0.42, [("M0", 8.2, True, C.ink)], C.gray, C.gray_line, 1.0)
    text(slide, 1.57, 1.84, 0.92, 0.20, [("Pair construction", 5.7, True, C.ink), ("SemanticSTFTemporalDataset", 4.6, False, C.dim)], PP_ALIGN.CENTER)

    stack_block(slide, 2.63, 1.28, 0.53, 0.56, "B1", "PT-v3m1", C.orange, C.orange_dark)
    box(slide, 3.42, 1.32, 0.59, 0.48, [("TRF", 8.0, True, C.ink), ("Temporal ref.", 5.2, False, C.dim)], C.blue, C.blue_dark, 1.0)
    stack_block(slide, 4.25, 1.28, 0.53, 0.56, "B2", "shared", C.orange, C.orange_dark)
    box(slide, 5.04, 1.32, 0.59, 0.48, [("TRF", 8.0, True, C.ink), ("Fusion ref.", 5.2, False, C.dim)], C.blue, C.blue_dark, 1.0)
    stack_block(slide, 5.86, 1.28, 0.53, 0.56, "B3", "shared", C.orange, C.orange_dark)
    stack_block(slide, 6.62, 1.28, 0.53, 0.56, "B4", "shared", C.orange, C.orange_dark)
    box(slide, 7.50, 1.31, 0.60, 0.50, [("F1", 8.0, True, C.ink), ("Global", 5.3, False, C.dim)], C.blue_light, C.blue_dark, 1.0)
    box(slide, 8.26, 1.31, 0.60, 0.50, [("F2", 8.0, True, C.ink), ("Local", 5.3, False, C.dim)], C.blue_light, C.blue_dark, 1.0)
    box(slide, 9.08, 1.35, 0.64, 0.42, [("Seg.", 7.6, True, C.ink), ("Linear(C,19)", 5.0, False, C.dim)], C.green_light, C.green_dark, 1.0)
    semantic_output(slide, 10.38, 1.18)

    # Main feature stream.
    arrow(slide, 1.52, 1.18, 1.78, 1.47)
    arrow(slide, 1.52, 1.80, 1.78, 1.58)
    stream_y = 1.55
    for x0, x1 in [
        (2.28, 2.63), (3.16, 3.42), (4.01, 4.25), (4.78, 5.04),
        (5.63, 5.86), (6.39, 6.62), (7.15, 7.50), (8.10, 8.26),
        (8.86, 9.08), (9.72, 10.38),
    ]:
        arrow(slide, x0, stream_y, x1, stream_y, C.line, 0.85)
    text(slide, 3.30, 1.06, 1.20, 0.12, [("shared t / t-1", 5.3, False, C.dim)], PP_ALIGN.CENTER)
    for x, label in [(1.72, "d=4"), (2.58, "c=64"), (4.20, "c=64"), (5.82, "c=128"), (6.58, "c=256"), (9.05, "d=19")]:
        text(slide, x, 1.12, 0.52, 0.12, [(label, 5.1, False, C.dim)], PP_ALIGN.CENTER)

    # Fewer, lighter top cross-layer lines.
    skip(slide, 2.88, 7.80, 0.60, 1.27)
    skip(slide, 4.50, 8.56, 0.45, 1.27)

    # Temporal reference generation beneath the trunk.
    box(slide, 2.55, 2.25, 3.02, 0.61, [], C.gray, C.gray_line, 0.85)
    text(slide, 2.85, 2.33, 2.42, 0.15, [("Temporal reference generation", 7.2, True, C.ink)], PP_ALIGN.CENTER)
    box(slide, 2.86, 2.58, 1.15, 0.20, [("Global context pooling", 5.2, False, C.ink)], C.white, C.blue_dark, 0.7)
    box(slide, 4.16, 2.58, 1.20, 0.20, [("Local grid correspondence", 5.2, False, C.ink)], C.white, C.blue_dark, 0.7)
    arrow(slide, 2.90, 1.84, 2.90, 2.25)
    arrow(slide, 4.50, 1.84, 4.50, 2.25)
    arrow(slide, 3.48, 2.25, 3.70, 1.80)
    arrow(slide, 5.02, 2.25, 5.33, 1.80)
    text(slide, 2.80, 2.91, 1.22, 0.16, [("g_tm1_point", 5.6, False, C.dim)], PP_ALIGN.CENTER)
    text(slide, 4.14, 2.91, 1.35, 0.16, [("l_tm1_point, valid_match", 5.6, False, C.dim)], PP_ALIGN.CENTER)

    box(slide, 0.82, 2.77, 1.40, 0.32, [("Mini-batch sampling", 6.1, True, C.ink), ("WeatherWeightedSampler", 4.7, False, C.dim)], C.purple_light, C.purple_dark, 0.9)
    arrow(slide, 1.52, 2.77, 1.98, 1.76, C.feedback, 0.65)

    # VFB and AWS moved upward and tightened.
    text(slide, 0.45, 3.55, 2.15, 0.16, [("VFB: Per-weather validation feedback", 6.8, True, C.ink)])
    vfb_y = 3.88
    box(slide, 0.48, vfb_y, 0.92, 0.38, [("Validation set", 5.5, True, C.ink), ("SemanticSTF val", 4.4, False, C.dim)], C.blue_light, C.blue_dark, 0.9)
    box(slide, 1.70, vfb_y, 0.92, 0.38, [("Weather split", 5.5, True, C.ink), ("_build_val_weather_map()", 4.0, False, C.dim)], C.blue_light, C.blue_dark, 0.9)
    box(slide, 2.92, vfb_y, 1.02, 0.38, [("Weather-wise stats", 5.4, True, C.ink), ("val mIoU / count", 4.4, False, C.dim)], C.blue_light, C.blue_dark, 0.9)
    box(slide, 4.24, vfb_y, 0.98, 0.38, [("Domain difficulty", 5.4, True, C.ink), ("d = 1 - mIoU", 4.4, False, C.dim)], C.purple_light, C.purple_dark, 0.9)
    for x0, x1 in [(1.40, 1.70), (2.62, 2.92), (3.94, 4.24)]:
        arrow(slide, x0, vfb_y + 0.19, x1, vfb_y + 0.19, C.line, 0.75)

    # Lighter validation feedback line, shorter and unobtrusive.
    arrow(slide, 10.98, 1.86, 10.98, 3.34, C.val, 0.45, dash=True, head=False)
    arrow(slide, 10.98, 3.34, 0.92, 3.34, C.val, 0.45, dash=True, head=False)
    arrow(slide, 0.92, 3.34, 0.92, vfb_y, C.val, 0.45, dash=True)

    text(slide, 6.22, 3.55, 2.15, 0.16, [("AWS: Adaptive weather curriculum sampler", 6.8, True, C.ink)])
    aws_y = 3.86
    aws = [
        (6.16, aws_y, 0.68, 0.34, "EMA", "mu=0.8", C.purple_light),
        (7.18, aws_y, 0.92, 0.34, "Difficulty dist.", "softmax(D/tau)", C.purple_light),
        (8.48, aws_y, 0.84, 0.34, "Curriculum prior", "p_base", C.gray),
        (6.52, 4.55, 1.14, 0.34, "Conservative update", "(1-beta)p_base + beta q", C.purple_light),
        (8.04, 4.55, 1.02, 0.34, "Bounded projection", "p_min <= p <= p_max", C.purple_light),
        (9.48, 4.46, 1.62, 0.50, "Next-stage sampling", "p_next -> sample weights", C.purple),
    ]
    for x, y, w, h, label, sub, fill in aws:
        box(slide, x, y, w, h, [(label, 5.4, True, C.ink), (sub, 4.2, False, C.dim)], fill, C.purple_dark if fill != C.gray else C.gray_line, 0.9)
    arrow(slide, 5.22, vfb_y + 0.19, 6.16, aws_y + 0.17, C.purple_dark, 0.7)
    arrow(slide, 6.84, aws_y + 0.17, 7.18, aws_y + 0.17)
    arrow(slide, 8.10, aws_y + 0.17, 8.48, aws_y + 0.17)
    arrow(slide, 7.63, aws_y + 0.34, 7.08, 4.55)
    arrow(slide, 8.90, aws_y + 0.34, 7.62, 4.55)
    arrow(slide, 7.66, 4.72, 8.04, 4.72)
    arrow(slide, 9.06, 4.72, 9.48, 4.72)

    # Short curriculum feedback loop, no outer-page border.
    arrow(slide, 10.25, 4.96, 10.25, 5.24, C.feedback, 0.55, dash=True, head=False)
    arrow(slide, 10.25, 5.24, 1.24, 5.24, C.feedback, 0.55, dash=True, head=False)
    arrow(slide, 1.24, 5.24, 1.24, 3.09, C.feedback, 0.55, dash=True)

    legend(slide)
    text(
        slide,
        0.45,
        6.98,
        12.42,
        0.24,
        [("Fig. 2. Overall framework of the proposed AWCT-TempSeg.", 10.5, False, C.ink, CAPTION_FONT)],
        PP_ALIGN.CENTER,
    )


def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--output", default=str(Path.home() / "Desktop" / "AWCT_TempSeg_Fig2_final_sci_style.pptx"))
    return p.parse_args()


def main():
    args = parse_args()
    prs = Presentation()
    prs.slide_width = I(SLIDE_W)
    prs.slide_height = I(SLIDE_H)
    build(prs)
    out = Path(args.output).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote refined SCI-style editable PPTX to {out}")


if __name__ == "__main__":
    main()
