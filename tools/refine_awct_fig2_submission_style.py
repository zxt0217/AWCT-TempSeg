#!/usr/bin/env python3
"""Submission-style SCI refinement for the AWCT-TempSeg framework figure.

This keeps the current method logic unchanged and only adjusts layout density,
typography, line weights, feedback routes, legend scale, and caption spacing.
All diagram elements are editable PowerPoint shapes/connectors/text boxes.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

from refine_awct_fig2_sci_style import (
    C,
    CAPTION_FONT,
    I,
    arrow,
    box,
    dot,
    lidar_frame,
    rgb,
    semantic_output,
    skip,
    text,
)


SLIDE_W = 13.333333
SLIDE_H = 7.5


def tiny_legend(slide):
    x, y = 8.10, 4.10
    text(slide, x, y, 0.42, 0.12, [("Legend", 5.3, True, C.ink)])
    items = [
        ("B", "Shared backbone block", C.orange),
        ("TRF", "Temporal reference/fusion module", C.blue),
        ("VFB", "Per-weather validation feedback", C.blue_light),
        ("AWS", "Adaptive weather sampler", C.purple),
    ]
    for i, (code, desc, fill) in enumerate(items):
        yy = y + 0.165 * (i + 1)
        box(slide, x + 0.48, yy, 0.27, 0.11, [(code, 4.3, True, C.ink)], fill, C.gray_line, 0.45)
        text(slide, x + 0.80, yy - 0.004, 1.55, 0.11, [(desc, 4.4, False, C.dim)])
    arrow(slide, x + 2.62, y + 0.20, x + 2.98, y + 0.20, C.line, 0.65)
    text(slide, x + 3.06, y + 0.145, 0.90, 0.12, [("Feature stream", 4.4, False, C.dim)])
    arrow(slide, x + 2.62, y + 0.46, x + 2.98, y + 0.46, "D4A06B", 0.55, dash=True)
    text(slide, x + 3.06, y + 0.405, 1.18, 0.12, [("Curriculum feedback", 4.4, False, C.dim)])


def build_compact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C.white)

    # No in-figure title. The caption carries the figure title, as in papers.

    # Top trunk, shifted upward and slightly tightened.
    lidar_frame(slide, 0.38, 0.72, "Current frame t")
    lidar_frame(slide, 0.38, 1.28, "Historical frame t-1")
    text(slide, 0.36, 0.53, 1.18, 0.13, [("Temporal LiDAR pair", 6.4, True, C.ink)], PP_ALIGN.CENTER)

    box(slide, 1.72, 1.05, 0.48, 0.40, [("M0", 8.0, True, C.ink)], C.gray, C.gray_line, 1.0)
    text(slide, 1.50, 1.52, 0.92, 0.17, [("Pair construction", 5.5, True, C.ink), ("SemanticSTFTemporalDataset", 4.3, False, C.dim)], PP_ALIGN.CENTER)

    box(slide, 2.52, 1.00, 0.51, 0.53, [("B1", 8.0, True, C.ink), ("PT-v3m1", 5.2, False, C.dim)], C.orange, C.orange_dark, 0.95)
    box(slide, 3.26, 1.04, 0.57, 0.46, [("TRF", 7.8, True, C.ink), ("Temporal ref.", 5.0, False, C.dim)], C.blue, C.blue_dark, 1.0)
    box(slide, 4.04, 1.00, 0.51, 0.53, [("B2", 8.0, True, C.ink), ("shared", 5.2, False, C.dim)], C.orange, C.orange_dark, 0.95)
    box(slide, 4.78, 1.04, 0.57, 0.46, [("TRF", 7.8, True, C.ink), ("Fusion ref.", 5.0, False, C.dim)], C.blue, C.blue_dark, 1.0)
    box(slide, 5.56, 1.00, 0.51, 0.53, [("B3", 8.0, True, C.ink), ("shared", 5.2, False, C.dim)], C.orange, C.orange_dark, 0.95)
    box(slide, 6.28, 1.00, 0.51, 0.53, [("B4", 8.0, True, C.ink), ("shared", 5.2, False, C.dim)], C.orange, C.orange_dark, 0.95)
    box(slide, 7.10, 1.03, 0.58, 0.48, [("F1", 7.8, True, C.ink), ("Global", 5.0, False, C.dim)], C.blue_light, C.blue_dark, 1.0)
    box(slide, 7.82, 1.03, 0.58, 0.48, [("F2", 7.8, True, C.ink), ("Local", 5.0, False, C.dim)], C.blue_light, C.blue_dark, 1.0)
    box(slide, 8.60, 1.07, 0.60, 0.40, [("Seg.", 7.4, True, C.ink), ("Linear(C,19)", 4.8, False, C.dim)], C.green_light, C.green_dark, 1.0)
    semantic_output(slide, 9.88, 0.89)

    # Main feature stream.
    arrow(slide, 1.50, 0.92, 1.72, 1.17)
    arrow(slide, 1.50, 1.50, 1.72, 1.30)
    stream_y = 1.25
    for x0, x1 in [
        (2.20, 2.52), (3.03, 3.26), (3.83, 4.04), (4.55, 4.78),
        (5.35, 5.56), (6.07, 6.28), (6.79, 7.10), (7.68, 7.82),
        (8.40, 8.60), (9.20, 9.88),
    ]:
        arrow(slide, x0, stream_y, x1, stream_y, C.line, 0.9)
    text(slide, 3.05, 0.82, 1.08, 0.11, [("shared t / t-1", 5.0, False, C.dim)], PP_ALIGN.CENTER)
    for x, label in [(1.66, "d=4"), (2.47, "c=64"), (3.99, "c=64"), (5.52, "c=128"), (6.24, "c=256"), (8.58, "d=19")]:
        text(slide, x, 0.86, 0.48, 0.11, [(label, 4.9, False, C.dim)], PP_ALIGN.CENTER)

    # Reduced, pale cross-layer links.
    skip(slide, 2.78, 7.38, 0.66, 0.99)
    skip(slide, 4.28, 8.10, 0.58, 0.99)

    # Temporal reference generation, light and close to the trunk.
    box(slide, 2.42, 1.82, 2.88, 0.52, [], "FAFAFA", C.gray_line, 0.65)
    text(slide, 2.68, 1.89, 2.34, 0.13, [("Temporal reference generation", 6.6, True, C.ink)], PP_ALIGN.CENTER)
    box(slide, 2.72, 2.11, 1.08, 0.17, [("Global context pooling", 4.7, False, C.ink)], C.white, C.blue_dark, 0.65)
    box(slide, 3.96, 2.11, 1.13, 0.17, [("Local grid correspondence", 4.7, False, C.ink)], C.white, C.blue_dark, 0.65)
    arrow(slide, 2.78, 1.53, 2.78, 1.82, C.line, 0.75)
    arrow(slide, 4.28, 1.53, 4.28, 1.82, C.line, 0.75)
    arrow(slide, 3.32, 1.82, 3.53, 1.50, C.line, 0.75)
    arrow(slide, 4.78, 1.82, 5.05, 1.50, C.line, 0.75)
    text(slide, 2.66, 2.36, 1.14, 0.13, [("g_tm1_point", 5.0, False, C.dim)], PP_ALIGN.CENTER)
    text(slide, 3.94, 2.36, 1.30, 0.13, [("l_tm1_point, valid_match", 5.0, False, C.dim)], PP_ALIGN.CENTER)

    box(slide, 8.10, 2.45, 1.30, 0.28, [("Mini-batch sampling", 5.6, True, C.ink), ("WeatherWeightedSampler", 4.3, False, C.dim)], C.purple_light, C.purple_dark, 0.9)

    # VFB/AWS mechanisms, raised and tightened.
    text(slide, 1.60, 2.83, 2.15, 0.14, [("VFB: Per-weather validation feedback", 6.6, True, C.ink)])
    vfb_y = 3.06
    box(slide, 1.55, vfb_y, 0.88, 0.33, [("Validation set", 5.0, True, C.ink), ("SemanticSTF val", 4.0, False, C.dim)], C.blue_light, C.blue_dark, 0.9)
    box(slide, 2.62, vfb_y, 0.88, 0.33, [("Weather split", 5.0, True, C.ink), ("_build_val_weather_map()", 3.6, False, C.dim)], C.blue_light, C.blue_dark, 0.9)
    box(slide, 3.70, vfb_y, 0.98, 0.33, [("Weather-wise stats", 4.9, True, C.ink), ("val mIoU / count", 3.9, False, C.dim)], C.blue_light, C.blue_dark, 0.9)
    box(slide, 4.90, vfb_y, 0.94, 0.33, [("Domain difficulty", 4.9, True, C.ink), ("d = 1 - mIoU", 3.9, False, C.dim)], C.purple_light, C.purple_dark, 0.9)
    for x0, x1 in [(2.43, 2.62), (3.50, 3.70), (4.68, 4.90)]:
        arrow(slide, x0, vfb_y + 0.165, x1, vfb_y + 0.165, C.line, 0.75)

    # Pale, short-weight validation feedback line to the validation statistics.
    arrow(slide, 10.48, 1.57, 10.48, 2.22, "BFDDB8", 0.55, dash=True, head=False)
    arrow(slide, 10.48, 2.22, 4.20, 2.96, "BFDDB8", 0.55, dash=True)

    text(slide, 6.18, 2.83, 2.20, 0.14, [("AWS: Adaptive weather curriculum sampler", 6.6, True, C.ink)])
    aws_y = 3.03
    aws = [
        (6.18, aws_y, 0.65, 0.30, "EMA", "mu=0.8", C.purple_light),
        (7.05, aws_y, 0.86, 0.30, "Difficulty dist.", "softmax(D/tau)", C.purple_light),
        (8.15, aws_y, 0.78, 0.30, "Curriculum prior", "p_base", C.gray),
        (6.40, 3.56, 1.08, 0.30, "Conservative update", "(1-beta)p_base + beta q", C.purple_light),
        (7.73, 3.56, 0.98, 0.30, "Bounded projection", "p_min <= p <= p_max", C.purple_light),
        (9.02, 3.47, 1.48, 0.42, "Next-stage sampling", "p_next -> sample weights", C.purple),
    ]
    for x, y, w, h, label, sub, fill in aws:
        box(slide, x, y, w, h, [(label, 4.9, True, C.ink), (sub, 3.8, False, C.dim)], fill, C.purple_dark if fill != C.gray else C.gray_line, 0.9)
    arrow(slide, 5.84, vfb_y + 0.165, 6.18, aws_y + 0.15, C.purple_dark, 0.75)
    arrow(slide, 6.83, aws_y + 0.15, 7.05, aws_y + 0.15, C.line, 0.75)
    arrow(slide, 7.91, aws_y + 0.15, 8.15, aws_y + 0.15, C.line, 0.75)
    arrow(slide, 7.45, aws_y + 0.30, 6.94, 3.56, C.line, 0.75)
    arrow(slide, 8.54, aws_y + 0.30, 7.48, 3.56, C.line, 0.75)
    arrow(slide, 7.48, 3.71, 7.73, 3.71, C.line, 0.75)
    arrow(slide, 8.71, 3.71, 9.02, 3.71, C.line, 0.75)

    # Short curriculum feedback loop: Next-stage sampling -> mini-batch sampler.
    arrow(slide, 9.76, 3.47, 9.76, 2.59, "D4A06B", 0.6, dash=True, head=False)
    arrow(slide, 9.76, 2.59, 9.40, 2.59, "D4A06B", 0.6, dash=True)

    tiny_legend(slide)
    text(
        slide,
        0.45,
        5.72,
        12.42,
        0.24,
        [("Fig. 2. Overall framework of the proposed AWCT-TempSeg.", 10.5, False, C.ink, CAPTION_FONT)],
        PP_ALIGN.CENTER,
    )


def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--output", default=str(Path.home() / "Desktop" / "AWCT_TempSeg_Fig2_final_submission_style.pptx"))
    return p.parse_args()


def main():
    args = parse_args()
    prs = Presentation()
    prs.slide_width = I(SLIDE_W)
    prs.slide_height = I(SLIDE_H)
    build_compact(prs)
    out = Path(args.output).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote submission-style editable PPTX to {out}")


if __name__ == "__main__":
    main()
