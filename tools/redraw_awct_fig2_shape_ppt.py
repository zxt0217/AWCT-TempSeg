#!/usr/bin/env python3
"""Redraw AWCT-TempSeg Fig. 2 as editable PowerPoint shapes.

This is a shape-level redraw in the visual language of a compact SCI method
figure: small rectangular network blocks, thin connectors, top skip links,
paper-style legend, and a caption. The slide is not a raster image.

Code-grounded anchors used before drawing:
- pointcept/datasets/semanticstf.py:
  SemanticSTFTemporalDataset, prev_* fields, has_prev, weather_map/_get_weather().
- pointcept/models/default.py:
  TempSegV2Segmentor, shared backbone, global/local temporal context,
  valid_match, temporal_fuse_global/local, seg_head.
- pointcept/engines/train.py:
  WeatherWeightedSampler, ValidationGuidedWeatherCurriculumTrainer,
  _build_val_weather_map(), validation mIoU, difficulty EMA, p_next weights.
- configs/semanticstf/semseg-pt-v3m1-0-tempseg-v2*.py:
  PT-v3m1 backbone, SemanticSTFTemporalDataset, tau/beta/mu/min/max/p_base.

Paper-expression blocks:
- TRF = Temporal reference/fusion module; groups global context and local
  same-grid correspondence from TempSegV2Segmentor.
- VFB = Per-weather validation feedback; groups weather split, val statistics,
  and difficulty computation.
- AWS = Adaptive weather sampler; groups EMA, conservative update, bounded
  projection, and WeatherWeightedSampler weight update.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


SLIDE_W = 13.333333
SLIDE_H = 7.5
FONT = "Arial"


class C:
    ink = "111111"
    dim = "4C4C4C"
    line = "2A2A2A"
    light_line = "6F6F6F"
    orange = "F7C9A8"
    orange_dark = "C95F16"
    orange_light = "FDE9D9"
    blue = "A9C7E8"
    blue_light = "DCEBF8"
    blue_dark = "2E75B6"
    green = "B9D9A6"
    green_light = "DDEFD5"
    green_dark = "548235"
    purple = "CFC3E6"
    purple_light = "ECE4F6"
    purple_dark = "6B57A6"
    gray = "F3F3F3"
    gray_line = "A6A6A6"
    white = "FFFFFF"
    feedback = "D66A1F"
    val = "6AA84F"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def I(v: float):
    return Inches(v)


def set_text(shape, lines, *, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0.025):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(0.015)
    tf.margin_bottom = I(0.015)
    for i, item in enumerate(lines):
        text, size, bold, color = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color or C.ink)


def shape_box(slide, x, y, w, h, lines, fill, line=C.line, lw=0.9, name=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x), I(y), I(w), I(h))
    if name:
        s.name = name
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line)
    s.line.width = Pt(lw)
    set_text(s, lines)
    return s


def text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, name=None):
    s = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    if name:
        s.name = name
    set_text(s, lines, align=align, margin=0)
    return s


def line_fmt(conn, color=C.line, lw=0.8, dash=False):
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(lw)
    if dash:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH


def add_arrowhead(connector):
    sp_pr = connector._element.spPr
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = OxmlElement("a:ln")
        sp_pr.append(ln)
    for old in ln.findall(qn("a:tailEnd")):
        ln.remove(old)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    ln.append(tail)


def arrow(slide, x1, y1, x2, y2, color=C.line, lw=0.8, dash=False, head=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    line_fmt(c, color, lw, dash)
    if head:
        add_arrowhead(c)
    return c


def dot(slide, x, y, color=C.line, size=0.025, lw=0.1):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(size), I(size))
    d.fill.solid()
    d.fill.fore_color.rgb = rgb(color)
    d.line.color.rgb = rgb(color)
    d.line.width = Pt(lw)
    return d


def block_stack(slide, x, y, w, h, label, sub, fill, line, layers=3):
    for i in range(layers - 1, 0, -1):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, I(x + 0.035 * i), I(y + 0.035 * i), I(w), I(h))
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(0.45)
    return shape_box(slide, x, y, w, h, [(label, 8.0, True, C.ink), (sub, 5.2, False, C.dim)], fill, line, 0.85)


def lidar_frame(slide, x, y, title, color):
    shape_box(slide, x, y, 1.22, 0.58, [], C.white, color, 0.75)
    text(slide, x + 0.08, y + 0.05, 1.04, 0.17, [(title, 6.4, True, C.ink)], PP_ALIGN.CENTER)
    pts = [
        (0.12, 0.28), (0.18, 0.38), (0.24, 0.25), (0.32, 0.43),
        (0.40, 0.32), (0.48, 0.24), (0.56, 0.44), (0.64, 0.31),
        (0.72, 0.39), (0.82, 0.27), (0.90, 0.46), (0.98, 0.35),
    ]
    for px, py in pts:
        dot(slide, x + px, y + py, color, 0.026)
    # A few scan-line strokes to suggest LiDAR sweep.
    for k in range(3):
        arrow(slide, x + 0.12, y + 0.48 - 0.06 * k, x + 1.05, y + 0.48 - 0.06 * k, color="B7B7B7", lw=0.35, head=False)


def semantic_output(slide, x, y):
    shape_box(slide, x, y, 1.34, 0.78, [], C.white, C.green_dark, 0.75)
    text(slide, x + 0.09, y + 0.04, 1.16, 0.22, [("Semantic prediction", 6.4, True, C.ink), ("seg_logits / labels", 5.1, False, C.dim)], PP_ALIGN.CENTER)
    colors = ["6AA84F", "3D85C6", "F1C232", "CC0000", "8E7CC3", "00A2E8", "FF99CC"]
    pts = [
        (0.15, 0.45, 0), (0.25, 0.62, 0), (0.39, 0.45, 1), (0.53, 0.65, 1),
        (0.67, 0.43, 2), (0.80, 0.58, 2), (0.47, 0.80, 3), (0.30, 0.78, 4),
        (0.72, 0.78, 4), (0.56, 0.48, 5), (0.90, 0.73, 6),
    ]
    for px, py, ci in pts:
        dot(slide, x + px * 1.34, y + py * 0.78, colors[ci], 0.035)


def skip_connection(slide, x0, x1, y_top, y_base):
    arrow(slide, x0, y_base, x0, y_top, C.light_line, 0.65, head=False)
    arrow(slide, x0, y_top, x1, y_top, C.light_line, 0.65, head=False)
    arrow(slide, x1, y_top, x1, y_base, C.light_line, 0.65, head=True)


def legend(slide):
    x, y = 6.05, 5.95
    text(slide, x, y - 0.05, 0.58, 0.18, [("Legend", 6.3, True, C.ink)])
    items = [
        ("B", "Shared backbone block", C.orange),
        ("TRF", "Temporal reference/fusion module", C.blue),
        ("VFB", "Per-weather validation feedback", C.blue_light),
        ("AWS", "Adaptive weather sampler", C.purple),
    ]
    for i, (abbr, desc, fill) in enumerate(items):
        yy = y + 0.27 * i
        shape_box(slide, x + 0.70, yy, 0.42, 0.18, [(abbr, 5.8, True, C.ink)], fill, C.gray_line, 0.45)
        text(slide, x + 1.18, yy - 0.01, 2.25, 0.18, [(desc, 5.8, False, C.dim)])
    arrow(slide, x + 3.30, y + 0.06, x + 3.78, y + 0.06, C.line, 0.8)
    text(slide, x + 3.86, y - 0.02, 1.15, 0.18, [("Feature stream", 5.8, False, C.dim)])
    arrow(slide, x + 3.30, y + 0.37, x + 3.78, y + 0.37, C.feedback, 0.8, dash=True)
    text(slide, x + 3.86, y + 0.29, 1.75, 0.18, [("Curriculum feedback", 5.8, False, C.dim)])


def build(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(C.white)

    text(slide, 0.45, 0.18, 3.8, 0.22, [("AWCT-TempSeg overall framework", 11.5, True, C.ink)])

    # Top compact network trunk, mapped from Fig. 2/Fig. 5 style modules.
    lidar_frame(slide, 0.38, 1.38, "Current frame t", C.orange_dark)
    lidar_frame(slide, 0.38, 2.08, "Historical frame t-1", C.orange_dark)
    text(slide, 0.38, 1.12, 1.28, 0.18, [("Temporal LiDAR pair", 7.0, True, C.ink)], PP_ALIGN.CENTER)

    shape_box(slide, 1.88, 1.80, 0.55, 0.46, [("M0", 8.3, True, C.ink)], C.gray, C.gray_line, 0.75)
    text(slide, 1.64, 2.34, 1.03, 0.24, [("Pair construction", 5.8, True, C.ink), ("SemanticSTFTemporalDataset", 4.8, False, C.dim)], PP_ALIGN.CENTER)

    b1 = block_stack(slide, 2.74, 1.74, 0.55, 0.58, "B1", "PT-v3m1", C.orange, C.orange_dark)
    trf1 = shape_box(slide, 3.55, 1.78, 0.62, 0.50, [("TRF", 8.0, True, C.ink), ("Temporal ref.", 5.0, False, C.dim)], C.blue, C.blue_dark, 0.8)
    b2 = block_stack(slide, 4.43, 1.74, 0.55, 0.58, "B2", "shared", C.orange, C.orange_dark)
    trf2 = shape_box(slide, 5.25, 1.78, 0.62, 0.50, [("TRF", 8.0, True, C.ink), ("Fusion ref.", 5.0, False, C.dim)], C.blue, C.blue_dark, 0.8)
    b3 = block_stack(slide, 6.08, 1.74, 0.55, 0.58, "B3", "shared", C.orange, C.orange_dark)
    b4 = block_stack(slide, 6.86, 1.74, 0.55, 0.58, "B4", "shared", C.orange, C.orange_dark)
    f1 = shape_box(slide, 7.75, 1.76, 0.63, 0.54, [("F1", 8.0, True, C.ink), ("Global", 5.2, False, C.dim)], C.blue_light, C.blue_dark, 0.75)
    f2 = shape_box(slide, 8.55, 1.76, 0.63, 0.54, [("F2", 8.0, True, C.ink), ("Local", 5.2, False, C.dim)], C.blue_light, C.blue_dark, 0.75)
    seg = shape_box(slide, 9.40, 1.80, 0.68, 0.46, [("Seg.", 7.4, True, C.ink), ("Linear(C,19)", 4.9, False, C.dim)], C.green_light, C.green_dark, 0.75)
    semantic_output(slide, 10.70, 1.64)

    # Main stream arrows.
    arrow(slide, 1.60, 1.67, 1.88, 1.95)
    arrow(slide, 1.60, 2.38, 1.88, 2.10)
    stream_y = 2.03
    for x0, x1 in [(2.43, 2.74), (3.29, 3.55), (4.17, 4.43), (4.98, 5.25), (5.87, 6.08), (6.63, 6.86), (7.41, 7.75), (8.38, 8.55), (9.18, 9.40), (10.08, 10.70)]:
        arrow(slide, x0, stream_y, x1, stream_y)
    arrow(slide, 7.41, stream_y, 7.75, stream_y)
    text(slide, 2.55, 1.48, 4.90, 0.16, [("shared weights for t and t-1; paired temporal information is reused across the trunk", 5.2, False, C.dim)], PP_ALIGN.CENTER)
    for x, label in [(1.88, "d=4"), (2.74, "c=64"), (4.43, "c=64"), (6.08, "c=128"), (6.86, "c=256"), (9.40, "d=19")]:
        text(slide, x, 1.55, 0.65, 0.15, [(label, 5.4, False, C.dim)], PP_ALIGN.CENTER)

    # Cross-layer links similar to the reference figure.
    skip_connection(slide, 2.98, 7.96, 0.92, 1.74)
    skip_connection(slide, 4.68, 8.86, 0.73, 1.74)
    skip_connection(slide, 6.33, 9.72, 0.54, 1.74)

    # Temporal reference generation, occupying the role of MPG in the reference.
    shape_box(slide, 2.68, 2.83, 3.10, 0.70, [], C.gray, C.gray_line, 0.65)
    text(slide, 3.02, 2.93, 2.42, 0.15, [("Temporal reference generation", 7.2, True, C.ink)], PP_ALIGN.CENTER)
    shape_box(slide, 2.92, 3.18, 1.20, 0.25, [("Global context pooling", 5.4, False, C.ink)], C.white, C.blue_dark, 0.45)
    shape_box(slide, 4.25, 3.18, 1.28, 0.25, [("Local grid correspondence", 5.4, False, C.ink)], C.white, C.blue_dark, 0.45)
    arrow(slide, 3.02, 2.32, 3.02, 2.83, head=True)
    arrow(slide, 4.70, 2.32, 4.70, 2.83, head=True)
    arrow(slide, 3.52, 2.83, 3.76, 2.28, head=True)
    arrow(slide, 5.25, 2.83, 5.56, 2.28, head=True)
    # Formula-like feature tokens below, drawn as editable text boxes.
    text(slide, 2.82, 3.56, 1.30, 0.21, [("g_tm1_point", 6.0, False, C.dim)], PP_ALIGN.CENTER)
    text(slide, 4.25, 3.56, 1.35, 0.21, [("l_tm1_point, valid_match", 6.0, False, C.dim)], PP_ALIGN.CENTER)

    # Mini-batch sampler receives feedback and feeds pair construction.
    shape_box(slide, 0.86, 3.16, 1.45, 0.34, [("Mini-batch sampling", 6.2, True, C.ink), ("WeatherWeightedSampler", 4.9, False, C.dim)], C.purple_light, C.purple_dark, 0.65)
    arrow(slide, 1.58, 3.16, 2.02, 2.26, C.feedback, 0.65)

    # Validation feedback branch; no enclosing large card.
    text(slide, 0.45, 4.23, 2.1, 0.18, [("VFB: Per-weather validation feedback", 7.0, True, C.ink)])
    vfb_y = 4.62
    shape_box(slide, 0.48, vfb_y, 0.95, 0.42, [("Validation set", 5.7, True, C.ink), ("SemanticSTF val", 4.6, False, C.dim)], C.blue_light, C.blue_dark, 0.65)
    shape_box(slide, 1.76, vfb_y, 0.95, 0.42, [("Weather split", 5.7, True, C.ink), ("_build_val_weather_map()", 4.3, False, C.dim)], C.blue_light, C.blue_dark, 0.65)
    shape_box(slide, 3.04, vfb_y, 1.05, 0.42, [("Weather-wise stats", 5.6, True, C.ink), ("val mIoU / count", 4.5, False, C.dim)], C.blue_light, C.blue_dark, 0.65)
    shape_box(slide, 4.42, vfb_y, 1.02, 0.42, [("Domain difficulty", 5.6, True, C.ink), ("d = 1 - mIoU", 4.6, False, C.dim)], C.purple_light, C.purple_dark, 0.65)
    for x0, x1 in [(1.43, 1.76), (2.71, 3.04), (4.09, 4.42)]:
        arrow(slide, x0, vfb_y + 0.21, x1, vfb_y + 0.21)

    # Dashed validation feedback from prediction to VFB.
    arrow(slide, 11.35, 2.42, 11.35, 3.96, C.val, 0.65, dash=True, head=False)
    arrow(slide, 11.35, 3.96, 0.95, 3.96, C.val, 0.65, dash=True, head=False)
    arrow(slide, 0.95, 3.96, 0.95, vfb_y, C.val, 0.65, dash=True, head=True)

    # Adaptive weather sampler mechanism, in a compact Fig. 6-like subfigure.
    text(slide, 6.25, 4.23, 2.0, 0.18, [("AWS: Adaptive weather curriculum sampler", 7.0, True, C.ink)])
    aws_y = 4.58
    modules = [
        (6.20, aws_y, 0.72, 0.38, "EMA", "mu=0.8", C.purple_light),
        (7.28, aws_y, 0.98, 0.38, "Difficulty dist.", "softmax(D/tau)", C.purple_light),
        (8.65, aws_y, 0.86, 0.38, "Curriculum prior", "p_base", C.gray),
        (6.60, 5.34, 1.22, 0.38, "Conservative update", "(1-beta)p_base + beta q", C.purple_light),
        (8.18, 5.34, 1.05, 0.38, "Bounded projection", "p_min <= p <= p_max", C.purple_light),
        (9.72, 5.25, 1.72, 0.52, "Next-stage sampling", "p_next -> sample weights", C.purple),
    ]
    for x, y, w, h, title, sub, fill in modules:
        shape_box(slide, x, y, w, h, [(title, 5.5, True, C.ink), (sub, 4.3, False, C.dim)], fill, C.purple_dark if fill != C.gray else C.gray_line, 0.65)
    arrow(slide, 5.44, vfb_y + 0.21, 6.20, aws_y + 0.19, C.purple_dark, 0.7)
    arrow(slide, 6.92, aws_y + 0.19, 7.28, aws_y + 0.19)
    arrow(slide, 8.26, aws_y + 0.19, 8.65, aws_y + 0.19)
    arrow(slide, 7.75, aws_y + 0.38, 7.20, 5.34)
    arrow(slide, 9.08, aws_y + 0.38, 7.74, 5.34)
    arrow(slide, 7.82, 5.53, 8.18, 5.53)
    arrow(slide, 9.23, 5.53, 9.72, 5.53)

    # Dashed curriculum loop to mini-batch sampler, routed below text.
    arrow(slide, 10.58, 5.77, 10.58, 7.00, C.feedback, 0.7, dash=True, head=False)
    arrow(slide, 10.58, 7.00, 0.58, 7.00, C.feedback, 0.7, dash=True, head=False)
    arrow(slide, 0.58, 7.00, 0.58, 3.33, C.feedback, 0.7, dash=True, head=False)
    arrow(slide, 0.58, 3.33, 0.86, 3.33, C.feedback, 0.7, dash=True, head=True)

    legend(slide)
    text(slide, 0.42, 7.18, 12.45, 0.22, [("Fig. 2. Overall framework of the proposed AWCT-TempSeg.", 9.5, True, C.ink)], PP_ALIGN.CENTER)


def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--output", default=str(Path.home() / "Desktop" / "AWCT_TempSeg_Fig2_shape_redraw.pptx"))
    return p.parse_args()


def main():
    args = parse_args()
    prs = Presentation()
    prs.slide_width = I(SLIDE_W)
    prs.slide_height = I(SLIDE_H)
    build(prs)
    out = Path(args.output).expanduser().resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"Wrote editable shape redraw to {out}")


if __name__ == "__main__":
    main()
